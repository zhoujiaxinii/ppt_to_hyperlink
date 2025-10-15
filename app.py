from flask import Flask, request, jsonify
from werkzeug.exceptions import BadRequest
import os
import requests
import zipfile
import tempfile
import re
import shutil
from datetime import datetime
from pptx import Presentation
from qcloud_cos import CosConfig, CosS3Client
import logging
import time
from functools import wraps
import threading
from urllib.parse import urlparse
import platform

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Tencent Cloud COS Configuration
# 请在环境变量中设置以下配置
COS_SECRET_ID = os.environ.get("COS_SECRET_ID")
COS_SECRET_KEY = os.environ.get("COS_SECRET_KEY")
COS_REGION = os.environ.get("COS_REGION")
COS_BUCKET = os.environ.get("COS_BUCKET")

# Configuration constants - Optimized for 1-minute performance
REQUEST_TIMEOUT = 15  # seconds - reduced for faster response
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB - reduced to ensure faster processing
MAX_RETRIES = 2  # reduced retries for faster failure
RETRY_DELAY = 0.5  # seconds - reduced delay
API_TIMEOUT = 60  # API总超时时间：1分钟
PROCESSING_TIMEOUT = 45  # 处理超时：45秒，为响应留余量

# Initialize COS client with error handling
cos_client = None
if COS_SECRET_ID and COS_SECRET_KEY and COS_REGION and COS_BUCKET:
    try:
        cos_config = CosConfig(Region=COS_REGION, SecretId=COS_SECRET_ID, SecretKey=COS_SECRET_KEY)
        cos_client = CosS3Client(cos_config)
        logger.info("COS client initialized successfully")
    except Exception as e:
        logger.error(f"Failed to initialize COS client: {str(e)}")
        cos_client = None
else:
    logger.warning("COS credentials not provided. COS functionality will be disabled.")

class TimeoutException(Exception):
    """Custom timeout exception"""
    pass

def with_timeout(seconds):
    """Decorator to add timeout to functions (Windows compatible)"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            result = [None]
            exception = [None]

            def target():
                try:
                    result[0] = func(*args, **kwargs)
                except Exception as e:
                    exception[0] = e

            thread = threading.Thread(target=target)
            thread.daemon = True
            thread.start()
            thread.join(seconds)

            if thread.is_alive():
                logger.warning(f"Operation timed out after {seconds} seconds")
                raise TimeoutException("Operation timed out")

            if exception[0]:
                raise exception[0]

            return result[0]
        return wrapper
    return decorator

def download_file_with_retry(url, filepath, max_retries=MAX_RETRIES):
    """Download file with retry mechanism and proper error handling"""
    for attempt in range(max_retries):
        try:
            logger.info(f"Downloading file (attempt {attempt + 1}/{max_retries})")

            # Validate URL
            parsed_url = urlparse(url)
            if not parsed_url.scheme or not parsed_url.netloc:
                raise ValueError(f"Invalid URL: {url}")

            # Make request with timeout and streaming - optimized for speed
            response = requests.get(
                url,
                stream=True,
                timeout=(5, REQUEST_TIMEOUT),  # (connect_timeout, read_timeout) - faster connect
                headers={
                    'User-Agent': 'PPT-Hyperlink-Converter/1.0',
                    'Accept-Encoding': 'gzip, deflate'  # Enable compression
                }
            )
            response.raise_for_status()

            # Check content length
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > MAX_FILE_SIZE:
                raise ValueError(f"File too large: {content_length} bytes (max: {MAX_FILE_SIZE})")

            # Download with size checking - optimized chunk size
            total_size = 0
            with open(filepath, 'wb') as f:
                for chunk in response.iter_content(chunk_size=65536):  # 64KB chunks for faster download
                    if chunk:
                        total_size += len(chunk)
                        if total_size > MAX_FILE_SIZE:
                            raise ValueError(f"File too large: {total_size} bytes (max: {MAX_FILE_SIZE})")
                        f.write(chunk)

            logger.info(f"File downloaded successfully: {total_size} bytes")
            return total_size

        except (requests.RequestException, ValueError) as e:
            logger.warning(f"Download attempt {attempt + 1} failed: {str(e)}")
            if attempt == max_retries - 1:
                raise
            time.sleep(RETRY_DELAY * (attempt + 1))  # Exponential backoff
        except Exception as e:
            logger.error(f"Unexpected error during download: {str(e)}")
            raise

@with_timeout(20)  # 20 seconds timeout for extraction
def extract_links_from_pptx(pptx_path):
    """
    Extract media and game links from PPTX file by examining its XML content

    Args:
        pptx_path (str): Path to the PPTX file

    Returns:
        set: Set of unique links found in the PPTX file
    """
    links = set()

    # Improved link patterns with better performance
    # Optimized patterns with pre-compilation - fixed to handle complex URLs
    media_extensions = 'mp3|mp4|wav|avi|mov|wmv|flv|ogg|webm'
    media_pattern = fr'https?://[^\s\)\(\（\）]+\.(?:{media_extensions})'

    # Game link pattern (optimized) - supports both formats:
    # 1. index.html?data_url=*.json
    # 2. index.html?data_url=*.json&studentId=*
    # Simplified to handle complex URLs with encoding
    game_pattern = r'https?://[^\s]+/index\.html\?data_url=https?://[^\s&]+\.json(?:&[^\s]*)?'

    # Pre-compile regex patterns for better performance
    media_regex = re.compile(media_pattern, re.IGNORECASE)
    game_regex = re.compile(game_pattern, re.IGNORECASE)

    try:
        # Extract PPTX file (which is essentially a ZIP file)
        with tempfile.TemporaryDirectory() as temp_dir:
            with zipfile.ZipFile(pptx_path, 'r') as zip_file:
                zip_file.extractall(temp_dir)

            # Walk through all files in the extracted directory
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    if file.endswith('.xml') or file.endswith('.rels'):
                        file_path = os.path.join(root, file)
                        try:
                            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                                content = f.read()

                                # Find media links (improved performance with compiled regex)
                                media_matches = media_regex.findall(content)
                                links.update(media_matches)

                                # Find game links (compiled regex)
                                game_matches = game_regex.findall(content)
                                links.update(game_matches)

                        except Exception as e:
                            logger.warning(f"Error reading file {file_path}: {str(e)}")
                            continue

    except Exception as e:
        logger.error(f"Error extracting links from PPTX: {str(e)}")
        raise

    # Clean up links by removing XML tags and entities
    cleaned_links = set()
    for link in links:
        # Remove XML tags and entities
        cleaned_link = re.sub(r'<[^>]*>', '', link)  # Remove XML tags
        cleaned_link = re.sub(r'&amp;', '&', cleaned_link)  # Replace &amp; with &
        cleaned_link = cleaned_link.strip()
        if cleaned_link.startswith('http'):
            cleaned_links.add(cleaned_link)

    return cleaned_links

def get_friendly_link_text(link):
    """
    Convert a link URL to friendly display text

    Args:
        link (str): The original link URL

    Returns:
        str: Friendly display text like "点击音频", "点击视频", "点击游戏"
    """
    link_lower = link.lower()

    # Check for game links first (more specific pattern)
    # Support both formats: index.html?data_url=*.json and index.html?data_url=*.json&studentId=*
    if 'index.html?data_url=' in link_lower and '.json' in link_lower:
        return "点击游戏"

    # Check for audio files
    audio_extensions = ['.mp3', '.wav', '.ogg']
    if any(ext in link_lower for ext in audio_extensions):
        return "点击音频"

    # Check for video files
    video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm']
    if any(ext in link_lower for ext in video_extensions):
        return "点击视频"

    # Default fallback - shouldn't happen with current regex patterns
    return "点击链接"

@with_timeout(15)  # 15 seconds timeout for hyperlink addition
def add_hyperlinks_to_pptx(pptx_path, links, output_path):
    """
    Add hyperlinks to text in PPTX file that matches the extracted links

    Args:
        pptx_path (str): Path to the input PPTX file
        links (set): Set of links to convert to hyperlinks
        output_path (str): Path for the output PPTX file
    """
    try:
        # Load the presentation
        prs = Presentation(pptx_path)

        # Track conversions for logging
        conversions_made = 0

        # Iterate through all slides
        for slide_idx, slide in enumerate(prs.slides):
            # Iterate through all shapes in the slide
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    # Iterate through all paragraphs in the text frame
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        paragraph_text = paragraph.text

                        # Improved link matching - check if paragraph contains link-like content
                        paragraph_lower = paragraph_text.lower()
                        matched_link = None

                        # Try exact match first
                        for link in links:
                            if link.lower() in paragraph_lower:
                                matched_link = link
                                break

                        # If no exact match, try partial matching for links
                        if not matched_link:
                            for link in links:
                                # Extract domain and key parts for matching
                                if 'index.html' in link and 'index.html' in paragraph_lower:
                                    # For game links, match by domain and index.html
                                    link_domain = link.split('/')[2] if '://' in link else ''
                                    if link_domain and link_domain in paragraph_lower:
                                        matched_link = link
                                        break
                                elif any(ext in link for ext in ['.mp4', '.mp3', '.wav']) and any(ext in paragraph_lower for ext in ['.mp4', '.mp3', '.wav']):
                                    # For media links, match by domain and extension
                                    link_domain = link.split('/')[2] if '://' in link else ''
                                    if link_domain and link_domain in paragraph_lower:
                                        matched_link = link
                                        break

                        if matched_link:
                            logger.info(f"Found link '{matched_link}' in slide {slide_idx + 1}, shape {shape_idx + 1}, paragraph {para_idx + 1}")

                            # Get friendly display text for the link
                            friendly_text = get_friendly_link_text(matched_link)

                            # Clear existing runs
                            paragraph.clear()

                            # Add new run with hyperlink using friendly text
                            run = paragraph.add_run()
                            run.text = friendly_text
                            run.hyperlink.address = matched_link

                            logger.info(f"Converted '{matched_link}' to display text '{friendly_text}'")
                            conversions_made += 1

        # Save the modified presentation
        prs.save(output_path)
        logger.info(f"Successfully processed PPTX file. Made {conversions_made} hyperlink conversions.")

    except Exception as e:
        logger.error(f"Error adding hyperlinks to PPTX: {str(e)}")
        raise

def validate_cos_config():
    """Validate COS configuration"""
    if not cos_client:
        raise RuntimeError("COS client not initialized")

    if not all([COS_SECRET_ID, COS_SECRET_KEY, COS_REGION, COS_BUCKET]):
        missing = []
        if not COS_SECRET_ID: missing.append('COS_SECRET_ID')
        if not COS_SECRET_KEY: missing.append('COS_SECRET_KEY')
        if not COS_REGION: missing.append('COS_REGION')
        if not COS_BUCKET: missing.append('COS_BUCKET')
        raise RuntimeError(f"Missing COS configuration: {', '.join(missing)}")

def upload_to_cos(file_path, cos_key, max_retries=MAX_RETRIES):
    """
    Upload file to Tencent Cloud COS

    Args:
        file_path (str): Local file path
        cos_key (str): Key (path) in COS bucket

    Returns:
        str: COS download URL
    """
    validate_cos_config()

    for attempt in range(max_retries):
        try:
            logger.info(f"Uploading to COS (attempt {attempt + 1}/{max_retries})")

            # Check file exists and size
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

            file_size = os.path.getsize(file_path)
            if file_size > MAX_FILE_SIZE:
                raise ValueError(f"File too large for upload: {file_size} bytes")

            # Upload file to COS
            with open(file_path, 'rb') as file_data:
                cos_client.put_object(
                    Bucket=COS_BUCKET,
                    Body=file_data,
                    Key=cos_key,
                    StorageClass='STANDARD'
                )

            # Construct download URL
            download_url = f"https://{COS_BUCKET}.cos.{COS_REGION}.myqcloud.com/{cos_key}"
            logger.info(f"File uploaded to COS successfully: {download_url}")

            return download_url

        except Exception as e:
            logger.warning(f"Upload attempt {attempt + 1} failed: {str(e)}")
            if attempt == max_retries - 1:
                logger.error(f"All upload attempts failed. Last error: {str(e)}")
                raise
            time.sleep(RETRY_DELAY * (attempt + 1))

@app.route('/process_pptx', methods=['POST'])
def process_pptx():
    """
    Process PPTX file: download, extract links, add hyperlinks, upload to COS

    Expected JSON payload:
    {
        "pptx_url": "https://example.com/file.pptx"
    }

    Returns:
        JSON response with COS download URL
    """
    try:
        # Get request data
        data = request.get_json()
        if not data or 'pptx_url' not in data:
            raise BadRequest("Missing 'pptx_url' in request body")

        pptx_url = data['pptx_url']
        start_time = time.time()
        logger.info(f"Processing PPTX from URL: {pptx_url}")
        logger.info(f"API timeout limit: {API_TIMEOUT} seconds")

        # Add timeout monitoring within the request context
        import threading
        import signal

        def timeout_handler():
            elapsed = time.time() - start_time
            if elapsed > API_TIMEOUT:
                logger.error(f"Request timeout after {API_TIMEOUT} seconds")
                # We can't directly return from here, but we can log it
                return True
            return False

        # Validate COS configuration first
        validate_cos_config()

        # Create temporary directory for processing
        with tempfile.TemporaryDirectory() as temp_dir:
            # Download PPTX file with retry mechanism
            logger.info("Downloading PPTX file...")
            input_pptx_path = os.path.join(temp_dir, "input.pptx")

            try:
                download_start = time.time()
                file_size = download_file_with_retry(pptx_url, input_pptx_path)
                download_time = time.time() - download_start
                logger.info(f"Downloaded PPTX file to: {input_pptx_path} ({file_size} bytes) in {download_time:.2f}s")

                # Check if we have enough time left
                elapsed_time = time.time() - start_time
                if elapsed_time > API_TIMEOUT * 0.3:  # If download took more than 30% of time
                    logger.warning(f"Download took {download_time:.2f}s, may not have enough time for processing")

            except Exception as e:
                logger.error(f"Failed to download file: {str(e)}")
                raise

            # Extract links from PPTX
            logger.info("Extracting links from PPTX...")
            extract_start = time.time()
            links = extract_links_from_pptx(input_pptx_path)
            extract_time = time.time() - extract_start
            logger.info(f"Found {len(links)} links in {extract_time:.2f}s: {list(links)}")

            if not links:
                logger.warning("No links found in PPTX file")
                return jsonify({
                    "success": False,
                    "message": "No media or game links found in the PPTX file"
                }), 400

            # Add hyperlinks to PPTX
            logger.info("Adding hyperlinks to PPTX...")
            hyperlink_start = time.time()
            output_pptx_path = os.path.join(temp_dir, "output.pptx")
            add_hyperlinks_to_pptx(input_pptx_path, links, output_pptx_path)
            hyperlink_time = time.time() - hyperlink_start
            logger.info(f"Added hyperlinks in {hyperlink_time:.2f}s")

            # Generate unique filename for COS
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            cos_key = f"processed_pptx/hyperlink_converted_{timestamp}.pptx"

            # Upload to COS
            logger.info("Uploading processed PPTX to COS...")
            upload_start = time.time()
            download_url = upload_to_cos(output_pptx_path, cos_key)
            upload_time = time.time() - upload_start
            total_time = time.time() - start_time
            logger.info(f"Uploaded to COS in {upload_time:.2f}s")
            logger.info(f"Total processing time: {total_time:.2f}s")

            return jsonify({
                "success": True,
                "message": "PPTX file processed successfully",
                "download_url": download_url,
                "links_found": list(links),
                "links_converted": len(links),
                "processing_time": round(total_time, 2),
                "performance": {
                    "download_time": round(download_time, 2),
                    "extract_time": round(extract_time, 2),
                    "hyperlink_time": round(hyperlink_time, 2),
                    "upload_time": round(upload_time, 2),
                    "total_time": round(total_time, 2)
                }
            })

    except (requests.RequestException, ValueError) as e:
        logger.error(f"Error downloading PPTX file: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"Error downloading PPTX file: {str(e)}",
            "error_type": "download_error"
        }), 400

    except TimeoutException as e:
        logger.error(f"Operation timed out: {str(e)}")
        return jsonify({
            "success": False,
            "message": "Operation timed out. Please try again with a smaller file.",
            "error_type": "timeout_error"
        }), 408

    except RuntimeError as e:
        logger.error(f"Configuration error: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"Server configuration error: {str(e)}",
            "error_type": "config_error"
        }), 500

    except Exception as e:
        logger.error(f"Unexpected error processing PPTX: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"Unexpected error: {str(e)}",
            "error_type": "server_error"
        }), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "service": "PPT Hyperlink Converter"})

@app.route('/', methods=['GET'])
def index():
    """API documentation"""
    return jsonify({
        "service": "PPT Hyperlink Converter API",
        "version": "1.0.0",
        "endpoints": {
            "POST /process_pptx": {
                "description": "Process PPTX file to add hyperlinks",
                "payload": {
                    "pptx_url": "URL of the PPTX file to process"
                },
                "response": {
                    "success": "boolean",
                    "message": "string",
                    "download_url": "string (COS download URL)",
                    "links_found": "array of strings",
                    "links_converted": "number"
                }
            },
            "GET /health": "Health check endpoint"
        }
    })

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)