from flask import Flask, request, jsonify
from werkzeug.exceptions import BadRequest
import os
import requests
import zipfile
import tempfile
import re
import shutil
from datetime import datetime
from pptx import Presentation
from qcloud_cos import CosConfig, CosS3Client
import logging
import time
from functools import wraps
import threading
from urllib.parse import urlparse
import platform

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Tencent Cloud COS Configuration
# 请在环境变量中设置以下配置
COS_SECRET_ID = os.environ.get("COS_SECRET_ID")
COS_SECRET_KEY = os.environ.get("COS_SECRET_KEY")
COS_REGION = os.environ.get("COS_REGION")
COS_BUCKET = os.environ.get("COS_BUCKET")

# Configuration constants - Optimized for 1-minute performance
REQUEST_TIMEOUT = 15  # seconds - reduced for faster response
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB - reduced to ensure faster processing
MAX_RETRIES = 2  # reduced retries for faster failure
RETRY_DELAY = 0.5  # seconds - reduced delay
API_TIMEOUT = 180  # API总超时时间：3分钟
PROCESSING_TIMEOUT = 45  # 处理超时：45秒，为响应留余量

# Initialize COS client with error handling
cos_client = None
if COS_SECRET_ID and COS_SECRET_KEY and COS_REGION and COS_BUCKET:
    try:
        cos_config = CosConfig(Region=COS_REGION, SecretId=COS_SECRET_ID, SecretKey=COS_SECRET_KEY, Timeout=60) # Increased timeout to 60 seconds
        cos_client = CosS3Client(cos_config)
        logger.info("COS client initialized successfully")
    except Exception as e:
        logger.error(f"Failed to initialize COS client: {str(e)}")
        cos_client = None
else:
    logger.warning("COS credentials not provided. COS functionality will be disabled.")

class TimeoutException(Exception):
    """Custom timeout exception"""
    pass

def with_timeout(seconds):
    """Decorator to add timeout to functions (Windows compatible)"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            result = [None]
            exception = [None]

            def target():
                try:
                    result[0] = func(*args, **kwargs)
                except Exception as e:
                    exception[0] = e

            thread = threading.Thread(target=target)
            thread.daemon = True
            thread.start()
            thread.join(seconds)

            if thread.is_alive():
                logger.warning(f"Operation timed out after {seconds} seconds")
                raise TimeoutException("Operation timed out")

            if exception[0]:
                raise exception[0]

            return result[0]
        return wrapper
    return decorator

def download_file_with_retry(url, filepath, max_retries=MAX_RETRIES):
    """Download file with retry mechanism and proper error handling"""
    for attempt in range(max_retries):
        try:
            logger.info(f"Downloading file (attempt {attempt + 1}/{max_retries})")

            # Validate URL
            parsed_url = urlparse(url)
            if not parsed_url.scheme or not parsed_url.netloc:
                raise ValueError(f"Invalid URL: {url}")

            # Make request with timeout and streaming - optimized for speed
            response = requests.get(
                url,
                stream=True,
                timeout=(5, REQUEST_TIMEOUT),  # (connect_timeout, read_timeout) - faster connect
                headers={
                    'User-Agent': 'PPT-Hyperlink-Converter/1.0',
                    'Accept-Encoding': 'gzip, deflate'  # Enable compression
                }
            )
            response.raise_for_status()

            # Check content length
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > MAX_FILE_SIZE:
                raise ValueError(f"File too large: {content_length} bytes (max: {MAX_FILE_SIZE})")

            # Download with size checking - optimized chunk size
            total_size = 0
            with open(filepath, 'wb') as f:
                for chunk in response.iter_content(chunk_size=65536):  # 64KB chunks for faster download
                    if chunk:
                        total_size += len(chunk)
                        if total_size > MAX_FILE_SIZE:
                            raise ValueError(f"File too large: {total_size} bytes (max: {MAX_FILE_SIZE})")
                        f.write(chunk)

            logger.info(f"File downloaded successfully: {total_size} bytes")
            return total_size

        except (requests.RequestException, ValueError) as e:
            logger.warning(f"Download attempt {attempt + 1} failed: {str(e)}")
            if attempt == max_retries - 1:
                raise
            time.sleep(RETRY_DELAY * (attempt + 1))  # Exponential backoff
        except Exception as e:
            logger.error(f"Unexpected error during download: {str(e)}")
            raise

@with_timeout(20)  # 20 seconds timeout for extraction
def extract_links_from_pptx(pptx_path):
    """
    Extract media and game links from PPTX file by examining its XML content
    
    This function properly extracts clean URLs from PPTX XML content by:
    1. Using more comprehensive URL patterns
    2. Cleaning extracted URLs from XML artifacts
    3. Handling URL encoding properly
    4. Filtering out image URLs

    Args:
        pptx_path (str): Path to the PPTX file

    Returns:
        set: Set of unique clean links found in the PPTX file
    """
    links = set()

    try:
        # Extract PPTX file (which is essentially a ZIP file)
        with tempfile.TemporaryDirectory() as temp_dir:
            with zipfile.ZipFile(pptx_path, 'r') as zip_file:
                zip_file.extractall(temp_dir)

            # Walk through all files in the extracted directory
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    if file.endswith('.xml') or file.endswith('.rels'):
                        file_path = os.path.join(root, file)
                        try:
                            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                                content = f.read()
                                
                                # Extract all URLs from the content using comprehensive pattern
                                # This pattern captures complete URLs including query parameters
                                url_pattern = r'https?://[^\s<>"\']+(?:\?[^\s<>"\']*)?(?:#[^\s<>"\']*)?'
                                found_urls = re.findall(url_pattern, content, re.IGNORECASE)
                                
                                for url in found_urls:
                                    # Clean the URL by removing XML artifacts and quotes
                                    clean_url = _clean_extracted_url(url)
                                    if clean_url and _is_valid_target_url(clean_url):
                                        links.add(clean_url)

                        except Exception as e:
                            logger.warning(f"Error reading file {file_path}: {str(e)}")
                            continue

            # Also check for hyperlinks in slide content using python-pptx
            try:
                prs = Presentation(pptx_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink.address:
                                        clean_url = _clean_extracted_url(run.hyperlink.address)
                                        if clean_url and _is_valid_target_url(clean_url):
                                            links.add(clean_url)
            except Exception as e:
                logger.warning(f"Error extracting hyperlinks using python-pptx: {str(e)}")

    except Exception as e:
        logger.error(f"Error extracting links from PPTX: {str(e)}")
        raise

    logger.info(f"Extracted {len(links)} clean URLs: {list(links)}")
    return links

def _clean_extracted_url(url):
    """
    Clean extracted URL by removing XML artifacts and normalizing
    
    Args:
        url (str): Raw URL extracted from XML
        
    Returns:
        str: Clean URL or None if invalid
    """
    if not url:
        return None
        
    # Remove common XML artifacts
    url = url.strip()
    url = url.rstrip('&quot;').rstrip('"').rstrip("'")
    url = url.lstrip('&quot;').lstrip('"').lstrip("'")
    
    # Remove XML entities
    url = url.replace('&amp;', '&')
    url = url.replace('&lt;', '<')
    url = url.replace('&gt;', '>')
    url = url.replace('&quot;', '"')
    
    # Remove trailing XML tags or artifacts
    url = re.sub(r'[<>].*$', '', url)
    
    # Validate basic URL structure
    if not re.match(r'^https?://', url, re.IGNORECASE):
        return None
        
    return url

def _is_valid_target_url(url):
    """
    Check if URL is a valid target for hyperlink conversion
    
    Args:
        url (str): URL to validate
        
    Returns:
        bool: True if URL should be converted to hyperlink
    """
    if not url:
        return False
        
    # Convert to lowercase for case-insensitive matching
    url_lower = url.lower()
    
    # Exclude image URLs
    image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp', '.ico']
    if any(ext in url_lower for ext in image_extensions):
        return False
        
    # Include game links (index.html with data_url parameter)
    if 'index.html' in url_lower and 'data_url=' in url_lower:
        return True
        
    # Include media links
    media_extensions = ['.mp3', '.mp4', '.wav', '.avi', '.mov', '.wmv', '.flv', '.ogg', '.webm']
    if any(ext in url_lower for ext in media_extensions):
        return True
        
    # Include other non-image URLs (general web links)
    # Exclude common static resource patterns
    static_patterns = ['/static/', '/assets/', '/images/', '/img/', '/css/', '/js/']
    if any(pattern in url_lower for pattern in static_patterns):
        return False
        
    return True

@with_timeout(15)  # 15 seconds timeout for hyperlink addition
def add_hyperlinks_to_pptx(pptx_path, links, output_path):
    """
    Add hyperlinks to text in PPTX file that matches the extracted links
    
    This improved function:
    1. Preserves original text formatting and content
    2. Only converts URLs that appear as plain text to hyperlinks
    3. Maintains existing hyperlinks and other content
    4. Uses more sophisticated text matching and replacement

    Args:
        pptx_path (str): Path to the input PPTX file
        links (set): Set of links to convert to hyperlinks
        output_path (str): Path for the output PPTX file
    """
    try:
        # Load the presentation
        prs = Presentation(pptx_path)

        # Track conversions for logging
        conversions_made = 0
        
        # Convert links to a list for easier processing
        links_list = list(links)

        # Iterate through all slides
        for slide_idx, slide in enumerate(prs.slides):
            # Iterate through all shapes in the slide
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    # Process each paragraph in the text frame
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        # Process the paragraph to add hyperlinks while preserving formatting
                        if _process_paragraph_for_hyperlinks(paragraph, links_list):
                            conversions_made += 1
                            logger.info(f"Added hyperlink in slide {slide_idx + 1}, shape {shape_idx + 1}, paragraph {para_idx + 1}")

        # Save the modified presentation
        prs.save(output_path)
        logger.info(f"Successfully processed PPTX file. Made {conversions_made} hyperlink conversions.")

    except Exception as e:
        logger.error(f"Error adding hyperlinks to PPTX: {str(e)}")
        raise

def _process_paragraph_for_hyperlinks(paragraph, links_list):
    """
    Process a paragraph to convert plain text URLs to hyperlinks while preserving formatting
    
    Args:
        paragraph: python-pptx paragraph object
        links_list (list): List of URLs to convert to hyperlinks
        
    Returns:
        bool: True if any conversions were made
    """
    conversions_made = False
    
    # Get the full paragraph text
    full_text = paragraph.text
    if not full_text:
        return False
    
    # Check if any of our target links appear in the paragraph text
    for link in links_list:
        if link in full_text:
            # Check if this URL is already a hyperlink
            if _is_url_already_hyperlinked(paragraph, link):
                continue
                
            # Find the position of the URL in the text
            url_start = full_text.find(link)
            if url_start == -1:
                continue
                
            url_end = url_start + len(link)
            
            # Split the paragraph into three parts: before URL, URL, after URL
            before_text = full_text[:url_start]
            url_text = full_text[url_start:url_end]
            after_text = full_text[url_end:]
            
            # Store original formatting from the first run
            original_font = None
            if paragraph.runs:
                original_font = paragraph.runs[0].font
            
            # Clear the paragraph
            paragraph.clear()
            
            # Add the text before the URL (if any)
            if before_text:
                run_before = paragraph.add_run()
                run_before.text = before_text
                if original_font:
                    _copy_font_properties(original_font, run_before.font)
            
            # Add the URL as a hyperlink
            run_url = paragraph.add_run()
            run_url.text = url_text
            run_url.hyperlink.address = link
            if original_font:
                _copy_font_properties(original_font, run_url.font)
            
            # Add the text after the URL (if any)
            if after_text:
                run_after = paragraph.add_run()
                run_after.text = after_text
                if original_font:
                    _copy_font_properties(original_font, run_after.font)
            
            conversions_made = True
            logger.info(f"Converted URL to hyperlink: {link}")
            break  # Only process one URL per paragraph to avoid conflicts
    
    return conversions_made

def _is_url_already_hyperlinked(paragraph, url):
    """
    Check if a URL is already a hyperlink in the paragraph
    
    Args:
        paragraph: python-pptx paragraph object
        url (str): URL to check
        
    Returns:
        bool: True if URL is already a hyperlink
    """
    for run in paragraph.runs:
        if hasattr(run, 'hyperlink') and run.hyperlink.address:
            if run.hyperlink.address == url or run.text == url:
                return True
    return False

def _copy_font_properties(source_font, target_font):
    """
    Copy font properties from source to target font
    
    Args:
        source_font: Source font object
        target_font: Target font object
    """
    try:
        if hasattr(source_font, 'name') and source_font.name:
            target_font.name = source_font.name
        if hasattr(source_font, 'size') and source_font.size:
            target_font.size = source_font.size
        if hasattr(source_font, 'bold') and source_font.bold is not None:
            target_font.bold = source_font.bold
        if hasattr(source_font, 'italic') and source_font.italic is not None:
            target_font.italic = source_font.italic
        if hasattr(source_font, 'underline') and source_font.underline is not None:
            target_font.underline = source_font.underline
        if hasattr(source_font, 'color') and source_font.color:
            target_font.color.rgb = source_font.color.rgb
    except Exception as e:
        logger.warning(f"Could not copy font property: {str(e)}")

def validate_cos_config():
    """Validate COS configuration"""
    if not cos_client:
        raise RuntimeError("COS client not initialized")

    if not all([COS_SECRET_ID, COS_SECRET_KEY, COS_REGION, COS_BUCKET]):
        missing = []
        if not COS_SECRET_ID: missing.append('COS_SECRET_ID')
        if not COS_SECRET_KEY: missing.append('COS_SECRET_KEY')
        if not COS_REGION: missing.append('COS_REGION')
        if not COS_BUCKET: missing.append('COS_BUCKET')
        raise RuntimeError(f"Missing COS configuration: {', '.join(missing)}")

def upload_to_cos(file_path, cos_key, max_retries=MAX_RETRIES):
    """
    Upload file to Tencent Cloud COS

    Args:
        file_path (str): Local file path
        cos_key (str): Key (path) in COS bucket

    Returns:
        str: COS download URL
    """
    validate_cos_config()

    for attempt in range(max_retries):
        try:
            logger.info(f"Uploading to COS (attempt {attempt + 1}/{max_retries})")

            # Check file exists and size
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")

            file_size = os.path.getsize(file_path)
            if file_size > MAX_FILE_SIZE:
                raise ValueError(f"File too large for upload: {file_size} bytes")

            # Upload file to COS
            with open(file_path, 'rb') as file_data:
                cos_client.put_object(
                    Bucket=COS_BUCKET,
                    Body=file_data,
                    Key=cos_key,
                    StorageClass='STANDARD'
                )

            # Construct download URL
            download_url = f"https://{COS_BUCKET}.cos.{COS_REGION}.myqcloud.com/{cos_key}"
            logger.info(f"File uploaded to COS successfully: {download_url}")

            return download_url

        except Exception as e:
            logger.warning(f"Upload attempt {attempt + 1} failed: {str(e)}")
            if attempt == max_retries - 1:
                logger.error(f"All upload attempts failed. Last error: {str(e)}")
                raise
            time.sleep(RETRY_DELAY * (attempt + 1))

@app.route('/process_pptx', methods=['POST'])
def process_pptx():
    """
    Process PPTX file: download, extract links, add hyperlinks, upload to COS

    Expected JSON payload:
    {
        "pptx_url": "https://example.com/file.pptx"
    }

    Returns:
        JSON response with COS download URL
    """
    try:
        # Get request data
        data = request.get_json()
        if not data or 'pptx_url' not in data:
            raise BadRequest("Missing 'pptx_url' in request body")

        pptx_url = data['pptx_url']
        start_time = time.time()
        logger.info(f"Processing PPTX from URL: {pptx_url}")

        # Create temporary directory for processing
        with tempfile.TemporaryDirectory() as temp_dir:
            # Download PPTX file with retry mechanism
            logger.info("Downloading PPTX file...")
            input_pptx_path = os.path.join(temp_dir, "input.pptx")

            download_start = time.time()
            file_size = download_file_with_retry(pptx_url, input_pptx_path)
            download_time = time.time() - download_start
            logger.info(f"Downloaded PPTX file: {file_size} bytes in {download_time:.2f}s")

            # Extract links from PPTX
            logger.info("Extracting links from PPTX...")
            extract_start = time.time()
            links = extract_links_from_pptx(input_pptx_path)
            extract_time = time.time() - extract_start
            logger.info(f"Found {len(links)} links in {extract_time:.2f}s")

            if not links:
                logger.warning("No links found in PPTX file")
                return jsonify({
                    "success": False,
                    "message": "No media or game links found in the PPTX file"
                }), 400

            # Add hyperlinks to PPTX
            logger.info("Adding hyperlinks to PPTX...")
            hyperlink_start = time.time()
            output_pptx_path = os.path.join(temp_dir, "output.pptx")
            add_hyperlinks_to_pptx(input_pptx_path, links, output_pptx_path)
            hyperlink_time = time.time() - hyperlink_start
            logger.info(f"Added hyperlinks in {hyperlink_time:.2f}s")

            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            unique_filename = f"hyperlink_converted_{timestamp}.pptx"
            final_output_path = os.path.join("/tmp", unique_filename)
            
            # Copy the processed file to tmp directory
            shutil.copy2(output_pptx_path, final_output_path)
            
            total_time = time.time() - start_time
            logger.info(f"Total processing time: {total_time:.2f}s")

            return jsonify({
                "success": True,
                "message": "PPTX file processed successfully",
                "links_found": list(links),
                "links_converted": len(links),
                "processing_time": round(total_time, 2),
                "performance": {
                    "download_time": round(download_time, 2),
                    "extract_time": round(extract_time, 2),
                    "hyperlink_time": round(hyperlink_time, 2),
                    "total_time": round(total_time, 2)
                },
                "processed_file_path": final_output_path
            })

    except (requests.RequestException, ValueError) as e:
        logger.error(f"Error processing PPTX file: {str(e)}")
        return jsonify({
            "success": False,
            "message": str(e)
        }), 500
    except BadRequest as e:
        logger.error(f"Bad request: {str(e)}")
        return jsonify({
            "success": False,
            "message": str(e)
        }), 400
    except TimeoutException as e:
        logger.error(f"Processing timed out: {str(e)}")
        return jsonify({
            "success": False,
            "message": "Processing timed out. Please try again or with a smaller file."
        }), 504
    except Exception as e:
        logger.error(f"An unexpected error occurred: {str(e)}", exc_info=True)
        return jsonify({
            "success": False,
            "message": "An unexpected error occurred during processing."
        }), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "service": "PPT Hyperlink Converter"})

@app.route('/', methods=['GET'])
def index():
    """API documentation"""
    return jsonify({
        "service": "PPT Hyperlink Converter API",
        "version": "1.0.0",
        "endpoints": {
            "POST /process_pptx": {
                "description": "Process PPTX file to add hyperlinks",
                "payload": {
                    "pptx_url": "URL of the PPTX file to process"
                },
                "response": {
                    "success": "boolean",
                    "message": "string",
                    "links_found": "array of strings",
                    "links_converted": "number"
                }
            },
            "GET /health": "Health check endpoint"
        }
    })

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)
