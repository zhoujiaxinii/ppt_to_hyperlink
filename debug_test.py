#!/usr/bin/env python3
"""
调试测试脚本 - 深入分析PPT中的文本内容和链接匹配
"""

import os
import sys
from pptx import Presentation
from app import extract_links_from_pptx, get_friendly_link_text

def debug_pptx_content(pptx_path):
    """
    深入分析PPTX文件内容，找出为什么链接没有被转换
    """
    if not os.path.exists(pptx_path):
        print(f"文件不存在: {pptx_path}")
        return

    print(f"调试文件: {pptx_path}")
    print("=" * 80)

    # 1. 提取链接
    print("1. 提取到的链接:")
    links = extract_links_from_pptx(pptx_path)
    for i, link in enumerate(links, 1):
        friendly_text = get_friendly_link_text(link)
        print(f"  {i}. {friendly_text}: {link[:60]}...")

    print(f"\n总共找到 {len(links)} 个链接\n")

    # 2. 分析PPT文本内容
    print("2. 分析PPT中的文本内容:")
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- 幻灯片 {slide_idx + 1} ---")

        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                print(f"  形状 {shape_idx + 1}:")

                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    para_text = paragraph.text
                    if para_text.strip():  # 只显示非空段落
                        print(f"    段落 {para_idx + 1}: '{para_text}'")

                        # 检查是否包含链接关键词
                        contains_http = 'http' in para_text.lower()
                        contains_index_html = 'index.html' in para_text.lower()
                        contains_mp4 = '.mp4' in para_text.lower()

                        if contains_http or contains_index_html or contains_mp4:
                            print(f"      -> 包含链接关键词: http={contains_http}, index.html={contains_index_html}, .mp4={contains_mp4}")

                        # 检查与提取的链接的匹配情况
                        for link in links:
                            if link.lower() in para_text.lower():
                                print(f"      -> 完全匹配链接: {get_friendly_link_text(link)}")
                            elif any(keyword in para_text.lower() for keyword in ['http', 'index.html', '.mp4', '.json']):
                                print(f"      -> 部分匹配可能: {para_text[:50]}...")

def main():
    """主调试函数"""
    print("PPT Hyperlink Converter - 调试测试")
    print("=" * 80)

    # 测试文件路径
    test_file = "ppt测试案例/汉语拼音声调初体验：一声 (19).pptx"

    debug_pptx_content(test_file)

if __name__ == "__main__":
    main()